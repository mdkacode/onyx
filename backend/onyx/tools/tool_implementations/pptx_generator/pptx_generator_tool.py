import io
import json
from typing import Any
from typing import cast

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from sqlalchemy.orm import Session
from typing_extensions import override

from onyx.chat.emitter import Emitter
from onyx.configs.constants import FileOrigin
from onyx.file_store.file_store import get_default_file_store
from onyx.file_store.utils import build_frontend_file_url
from onyx.server.query_and_chat.placement import Placement
from onyx.server.query_and_chat.streaming_models import GeneratedPptx
from onyx.server.query_and_chat.streaming_models import Packet
from onyx.server.query_and_chat.streaming_models import PptxGenerationFinal
from onyx.server.query_and_chat.streaming_models import PptxGenerationToolStart
from onyx.tools.interface import Tool
from onyx.tools.models import ToolResponse
from onyx.tools.tool_implementations.pptx_generator.models import (
    FinalPptxGenerationResponse,
)
from onyx.utils.logger import setup_logger

logger = setup_logger()

TITLE_FIELD = "title"
SLIDES_FIELD = "slides"

# Theme color definitions
THEME_COLORS: dict[str, dict[str, RGBColor]] = {
    "professional": {
        "title": RGBColor(0x1B, 0x3A, 0x5C),
        "body": RGBColor(0x00, 0x00, 0x00),
    },
    "modern": {
        "title": RGBColor(0x00, 0x80, 0x80),
        "body": RGBColor(0x33, 0x33, 0x33),
    },
    "minimal": {
        "title": RGBColor(0x00, 0x00, 0x00),
        "body": RGBColor(0x33, 0x33, 0x33),
    },
}

MAX_BULLETS_PER_SLIDE = 6


class PptxGeneratorTool(Tool[None]):
    NAME = "generate_presentation"
    DESCRIPTION = (
        "Generate a professional PowerPoint presentation (.pptx) from a structured outline. "
        "Use this when a user asks you to create a presentation, slides, or PPT."
    )
    DISPLAY_NAME = "Presentation Generator"

    def __init__(
        self,
        tool_id: int,
        emitter: Emitter,
    ) -> None:
        super().__init__(emitter=emitter)
        self._id = tool_id

    @property
    def id(self) -> int:
        return self._id

    @property
    def name(self) -> str:
        return self.NAME

    @property
    def description(self) -> str:
        return self.DESCRIPTION

    @property
    def display_name(self) -> str:
        return self.DISPLAY_NAME

    @override
    @classmethod
    def is_available(cls, db_session: Session) -> bool:
        """Always available -- python-pptx is a standard dependency."""
        return True

    def tool_definition(self) -> dict:
        return {
            "type": "function",
            "function": {
                "name": self.name,
                "description": "Generate a PowerPoint presentation from structured content",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "title": {
                            "type": "string",
                            "description": "The presentation title",
                        },
                        "subtitle": {
                            "type": "string",
                            "description": "Optional subtitle for the title slide",
                        },
                        "slides": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "title": {"type": "string"},
                                    "bullet_points": {
                                        "type": "array",
                                        "items": {"type": "string"},
                                    },
                                    "notes": {
                                        "type": "string",
                                        "description": "Speaker notes for this slide",
                                    },
                                },
                                "required": ["title", "bullet_points"],
                            },
                            "description": "Array of slide content objects",
                        },
                        "theme": {
                            "type": "string",
                            "enum": ["professional", "modern", "minimal"],
                            "description": "Visual theme for the presentation",
                        },
                    },
                    "required": ["title", "slides"],
                },
            },
        }

    def emit_start(self, placement: Placement) -> None:
        self.emitter.emit(
            Packet(
                placement=placement,
                obj=PptxGenerationToolStart(),
            )
        )

    @staticmethod
    def _split_long_content(
        bullet_points: list[str],
        max_per_slide: int = MAX_BULLETS_PER_SLIDE,
    ) -> list[list[str]]:
        """Split a long list of bullet points into chunks."""
        if len(bullet_points) <= max_per_slide:
            return [bullet_points]
        return [
            bullet_points[i : i + max_per_slide]
            for i in range(0, len(bullet_points), max_per_slide)
        ]

    @staticmethod
    def _apply_theme(
        prs: Presentation,
        theme: str,
    ) -> None:
        """Apply color scheme to all slides in the presentation."""
        colors = THEME_COLORS.get(theme, THEME_COLORS["professional"])
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Title shapes use index 0 in placeholders
                            if (
                                shape.placeholder_format is not None
                                and shape.placeholder_format.idx == 0
                            ):
                                run.font.color.rgb = colors["title"]
                            else:
                                run.font.color.rgb = colors["body"]

    def _add_content_slide(
        self,
        prs: Presentation,
        slide_data: dict[str, Any],
        theme: str,
    ) -> int:
        """Add content slide(s) for a single slide_data entry.

        Returns the number of slides added (may be >1 if bullet points are split).
        """
        bullet_points: list[str] = slide_data.get("bullet_points", [])
        slide_title: str = slide_data.get("title", "")
        notes: str | None = slide_data.get("notes")

        chunks = self._split_long_content(bullet_points)
        colors = THEME_COLORS.get(theme, THEME_COLORS["professional"])
        slides_added = 0

        for chunk_idx, chunk in enumerate(chunks):
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)

            # Set title
            title_shape = slide.shapes.title
            if title_shape is not None:
                title_shape.text = slide_title
                if len(chunks) > 1:
                    title_shape.text = f"{slide_title} ({chunk_idx + 1}/{len(chunks)})"
                for run in title_shape.text_frame.paragraphs[0].runs:
                    run.font.color.rgb = colors["title"]

            # Add bullet points to body placeholder
            body_placeholder = slide.placeholders[1]
            tf = body_placeholder.text_frame
            tf.clear()
            for i, point in enumerate(chunk):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = point
                p.font.size = Pt(18)
                for run in p.runs:
                    run.font.color.rgb = colors["body"]

            # Add speaker notes (only on the first chunk slide)
            if notes and chunk_idx == 0:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes

            slides_added += 1

        return slides_added

    def run(
        self,
        placement: Placement,
        override_kwargs: None = None,  # noqa: ARG002
        **llm_kwargs: Any,
    ) -> ToolResponse:
        title = cast(str, llm_kwargs.get(TITLE_FIELD, "Untitled Presentation"))
        subtitle = cast(str | None, llm_kwargs.get("subtitle"))
        slides_data = cast(list[dict[str, Any]], llm_kwargs.get(SLIDES_FIELD, []))
        theme = cast(str, llm_kwargs.get("theme", "professional"))

        if theme not in THEME_COLORS:
            theme = "professional"

        try:
            prs = Presentation()
            colors = THEME_COLORS[theme]

            # --- Title slide (layout 0) ---
            title_slide_layout = prs.slide_layouts[0]
            title_slide = prs.slides.add_slide(title_slide_layout)
            title_shape = title_slide.shapes.title
            if title_shape is not None:
                title_shape.text = title
                for run in title_shape.text_frame.paragraphs[0].runs:
                    run.font.color.rgb = colors["title"]

            if subtitle and title_slide.placeholders[1] is not None:
                subtitle_placeholder = title_slide.placeholders[1]
                subtitle_placeholder.text = subtitle
                for run in subtitle_placeholder.text_frame.paragraphs[0].runs:
                    run.font.color.rgb = colors["body"]

            # --- Content slides ---
            total_content_slides = 0
            for slide_data in slides_data:
                total_content_slides += self._add_content_slide(prs, slide_data, theme)

            num_slides = 1 + total_content_slides  # title slide + content slides

            # Save to BytesIO buffer
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)

            # Persist to file store
            file_store = get_default_file_store()
            file_id = file_store.save_file(
                content=buffer,
                display_name=f"{title}.pptx",
                file_origin=FileOrigin.OTHER,
                file_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )

            file_url = build_frontend_file_url(file_id)

            generated_pptx = GeneratedPptx(
                file_id=file_id,
                url=file_url,
                title=title,
                num_slides=num_slides,
            )

            # Emit final packet
            self.emitter.emit(
                Packet(
                    placement=placement,
                    obj=PptxGenerationFinal(presentation=generated_pptx),
                )
            )

            final_response = FinalPptxGenerationResponse(
                file_id=file_id,
                file_url=file_url,
                title=title,
                num_slides=num_slides,
            )

            llm_facing_response = json.dumps(
                {
                    "file_id": file_id,
                    "title": title,
                    "num_slides": num_slides,
                    "message": (
                        f"Generated a {num_slides}-slide presentation "
                        f"titled '{title}'. The file is ready for download."
                    ),
                }
            )

            return ToolResponse(
                rich_response=final_response,
                llm_facing_response=llm_facing_response,
            )

        except Exception:
            logger.exception("Error generating PowerPoint presentation")
            raise
